"""Multi-Format Enterprise Exporter (.docx, .pptx, .xlsx) with Embedded Visualizations.

Enables exporting rich executive reports, presentation slide decks, and formatted
Excel workbooks complete with embedded charts, KPI summaries, and structured tables.
"""

from __future__ import annotations

import datetime
import io
import logging
from typing import Any

import pandas as pd
from app.agent.chart_renderer import render_chart_to_png

logger = logging.getLogger(__name__)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 1. Word Document (.docx) Exporter with Embedded Charts & Tables
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def export_report_to_docx(
    title: str,
    report: dict[str, Any],
    charts: list[dict[str, Any]] | None = None,
    created_at: float | None = None,
) -> io.BytesIO:
    """Render a comprehensive executive report as a styled .docx document with embedded charts."""
    from docx import Document
    from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement, parse_xml
    from docx.oxml.ns import nsdecls, qn
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()

    # Page setup: Normal 1-inch margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    # Document Header & Title
    created_dt = datetime.datetime.fromtimestamp(created_at) if created_at else datetime.datetime.now()
    created_str = created_dt.strftime("%d/%m/%Y %H:%M")

    # Document Main Title
    title_p = doc.add_paragraph()
    title_run = title_p.add_run(title.upper())
    title_run.font.name = "Segoe UI"
    title_run.font.size = Pt(20)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(16, 185, 129)  # Emerald Accent

    # Subtitle / Metadata
    sub_p = doc.add_paragraph()
    sub_run = sub_p.add_run(f"Báo Cáo Phân Tích Dữ Liệu Tự Động • Ngày tạo: {created_str} • Hệ thống GraphSheet AI")
    sub_run.font.name = "Segoe UI"
    sub_run.font.size = Pt(9.5)
    sub_run.font.color.rgb = RGBColor(107, 114, 128)

    # Horizontal Divider Line
    p_div = doc.add_paragraph()
    p_div_border = parse_xml(r'<w:pBdr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                             r'<w:bottom w:val="single" w:sz="12" w:space="1" w:color="10B981"/>'
                             r'</w:pBdr>')
    p_div._p.get_or_add_pPr().append(p_div_border)

    # 1. Executive Summary Callout Box
    exec_summary = report.get("executive_summary", "")
    if exec_summary:
        h1 = doc.add_heading(level=1)
        h1_run = h1.add_run("1. Tóm Tắt Điều Hành (Executive Summary)")
        h1_run.font.name = "Segoe UI"
        h1_run.font.color.rgb = RGBColor(31, 41, 55)

        # Callout Table
        callout_table = doc.add_table(rows=1, cols=1)
        callout_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        cell = callout_table.cell(0, 0)
        cell.width = Inches(6.5)

        # Shading and border for callout box
        shading = parse_xml(r'<w:shd {} w:fill="F0FDF4"/>'.format(nsdecls('w')))
        borders = parse_xml(r'<w:tcBorders {}><w:left w:val="single" w:sz="24" w:space="0" w:color="10B981"/>'
                            r'<w:top w:val="none"/><w:right w:val="none"/><w:bottom w:val="none"/></w:tcBorders>'.format(nsdecls('w')))
        cell._tc.get_or_add_tcPr().append(shading)
        cell._tc.get_or_add_tcPr().append(borders)

        cp = cell.paragraphs[0]
        c_run = cp.add_run(exec_summary)
        c_run.font.name = "Segoe UI"
        c_run.font.size = Pt(10.5)
        c_run.font.color.rgb = RGBColor(17, 24, 39)
        doc.add_paragraph()

    # 2. Key Findings
    findings = report.get("key_findings") or []
    if findings:
        h2 = doc.add_heading(level=1)
        h2_run = h2.add_run("2. Các Phát Hiện Trọng Tâm (Key Findings)")
        h2_run.font.name = "Segoe UI"
        h2_run.font.color.rgb = RGBColor(31, 41, 55)

        for item in findings:
            fp = doc.add_paragraph(style="List Bullet")
            fr = fp.add_run(str(item))
            fr.font.name = "Segoe UI"
            fr.font.size = Pt(10)
            fr.font.color.rgb = RGBColor(55, 65, 81)

    # 3. Embedded Charts Section
    if charts:
        h_chart = doc.add_heading(level=1)
        h_chart_run = h_chart.add_run("3. Trực Quan Hoá Số Liệu & Biểu Đồ")
        h_chart_run.font.name = "Segoe UI"
        h_chart_run.font.color.rgb = RGBColor(31, 41, 55)

        for idx, chart in enumerate(charts):
            chart_title = chart.get("title") or f"Biểu đồ phân tích {idx+1}"
            cp = doc.add_paragraph()
            c_run = cp.add_run(f"• {chart_title}")
            c_run.font.name = "Segoe UI"
            c_run.font.size = Pt(11)
            c_run.font.bold = True
            c_run.font.color.rgb = RGBColor(37, 99, 235)

            # Render Chart to PNG image
            img_buf = render_chart_to_png(chart, width_in=6.2, height_in=3.6, dpi=220)
            if img_buf:
                img_p = doc.add_paragraph()
                img_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                img_p.add_run().add_picture(img_buf, width=Inches(6.0))
            
            # Optional data table under chart
            labels = chart.get("labels") or []
            values = chart.get("values") or []
            if labels and values and len(labels) == len(values) and len(labels) <= 10:
                table = doc.add_table(rows=len(labels) + 1, cols=2)
                table.alignment = WD_TABLE_ALIGNMENT.CENTER
                
                # Header row
                hdr_cells = table.rows[0].cells
                hdr_cells[0].text = chart.get("x_axis") or "Danh Mục"
                hdr_cells[1].text = chart.get("y_axis") or "Giá Trị"
                for hc in hdr_cells:
                    hc.paragraphs[0].runs[0].font.bold = True
                    hc.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                    shd = parse_xml(r'<w:shd {} w:fill="1F2937"/>'.format(nsdecls('w')))
                    hc._tc.get_or_add_tcPr().append(shd)

                # Data rows
                for r_idx, (lbl, val) in enumerate(zip(labels, values)):
                    row_cells = table.rows[r_idx + 1].cells
                    row_cells[0].text = str(lbl)
                    try:
                        row_cells[1].text = f"{float(val):,.2f}"
                    except (ValueError, TypeError):
                        row_cells[1].text = str(val)
                    row_cells[1].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT

                    # Alternate shading
                    if r_idx % 2 == 1:
                        for rc in row_cells:
                            shd = parse_xml(r'<w:shd {} w:fill="F9FAFB"/>'.format(nsdecls('w')))
                            rc._tc.get_or_add_tcPr().append(shd)

                doc.add_paragraph()

    # 4. Anomalies
    anomalies = report.get("anomalies") or []
    if anomalies:
        h3 = doc.add_heading(level=1)
        h3_run = h3.add_run("4. Cảnh Báo Bất Thường (Anomalies & Risks)")
        h3_run.font.name = "Segoe UI"
        h3_run.font.color.rgb = RGBColor(220, 38, 38)

        for item in anomalies:
            ap = doc.add_paragraph(style="List Bullet")
            ar = ap.add_run(f"⚠️ {item}")
            ar.font.name = "Segoe UI"
            ar.font.size = Pt(10)
            ar.font.color.rgb = RGBColor(185, 28, 28)

    # 5. Strategic Recommendations
    recs = report.get("recommendations") or []
    if recs:
        h4 = doc.add_heading(level=1)
        h4_run = h4.add_run("5. Đề Xuất Hành Động (Strategic Recommendations)")
        h4_run.font.name = "Segoe UI"
        h4_run.font.color.rgb = RGBColor(31, 41, 55)

        for idx, item in enumerate(recs):
            rp = doc.add_paragraph()
            rr_num = rp.add_run(f"Bước {idx+1}: ")
            rr_num.font.name = "Segoe UI"
            rr_num.font.bold = True
            rr_num.font.color.rgb = RGBColor(16, 185, 129)
            rr_text = rp.add_run(str(item))
            rr_text.font.name = "Segoe UI"
            rr_text.font.size = Pt(10)
            rr_text.font.color.rgb = RGBColor(55, 65, 81)

    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    return out


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 2. PowerPoint (.pptx) Presentation Exporter
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def export_report_to_pptx(
    title: str,
    report: dict[str, Any],
    charts: list[dict[str, Any]] | None = None,
    created_at: float | None = None,
) -> io.BytesIO:
    """Render an executive 16:9 presentation slide deck with embedded charts."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.warning("[exporter] python-pptx is not installed, fallback to empty buffer")
        return io.BytesIO()

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)     # #0F172A
    ACCENT_GREEN = RGBColor(16, 185, 129) # #10B981
    ACCENT_BLUE = RGBColor(37, 99, 235)  # #2563EB
    TEXT_LIGHT = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(148, 163, 184)
    TEXT_DARK = RGBColor(30, 41, 59)

    created_dt = datetime.datetime.fromtimestamp(created_at) if created_at else datetime.datetime.now()
    created_str = created_dt.strftime("%d/%m/%Y")

    # ── SLIDE 1: Title Slide ──────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5)) # Rectangle
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = BG_DARK
    bg1.line.fill.background()

    # Title Text Box
    txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p_badge = tf.paragraphs[0]
    p_badge.text = "EXECUTIVE INTELLIGENCE REPORT • GRAPHSHEET AI"
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_GREEN

    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT

    p_sub = tf.add_paragraph()
    p_sub.text = f"Báo cáo phân tích số liệu tự động • Ngày xuất bản: {created_str}"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = TEXT_MUTED

    # ── SLIDE 2: Executive Summary & Key Highlights ──────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Header
    tb_hdr = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_hdr = tb_hdr.text_frame
    p_h = tf_hdr.paragraphs[0]
    p_h.text = "TỔNG QUAN ĐIỀU HÀNH & KẾT QUẢ TRỌNG TÂM"
    p_h.font.size = Pt(22)
    p_h.font.bold = True
    p_h.font.color.rgb = ACCENT_BLUE

    # Executive Summary Card
    exec_summary = report.get("executive_summary", "")
    card1 = slide2.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(240, 253, 244)
    card1.line.color.rgb = ACCENT_GREEN
    card1.line.width = Pt(1.5)

    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    p_c1_title = tf_c1.paragraphs[0]
    p_c1_title.text = "📌 Tóm tắt cốt lõi:"
    p_c1_title.font.bold = True
    p_c1_title.font.size = Pt(13)
    p_c1_title.font.color.rgb = ACCENT_GREEN

    p_c1_body = tf_c1.add_paragraph()
    p_c1_body.text = exec_summary or "Đã tổng hợp số liệu phân tích."
    p_c1_body.font.size = Pt(13)
    p_c1_body.font.color.rgb = TEXT_DARK

    # Key Findings Bullets
    findings = report.get("key_findings") or []
    if findings:
        tb_find = slide2.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.2))
        tf_find = tb_find.text_frame
        tf_find.word_wrap = True
        p_f_head = tf_find.paragraphs[0]
        p_f_head.text = "🔍 Các phát hiện chính:"
        p_f_head.font.bold = True
        p_f_head.font.size = Pt(15)
        p_f_head.font.color.rgb = TEXT_DARK

        for f in findings[:4]:
            pf = tf_find.add_paragraph()
            pf.text = f"• {f}"
            pf.font.size = Pt(13)
            pf.font.color.rgb = RGBColor(51, 65, 85)

    # ── SLIDE 3..N: Chart Deep Dives ─────────────────────────────────────
    if charts:
        for idx, chart in enumerate(charts[:4]):
            slide_c = prs.slides.add_slide(blank_layout)
            
            # Slide Header
            c_title = chart.get("title") or f"Biểu đồ phân tích {idx+1}"
            tb_chdr = slide_c.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
            p_ch = tb_chdr.text_frame.paragraphs[0]
            p_ch.text = f"PHÂN TÍCH TRỰC QUAN: {c_title.upper()}"
            p_ch.font.size = Pt(20)
            p_ch.font.bold = True
            p_ch.font.color.rgb = ACCENT_BLUE

            # Render Chart Image
            img_buf = render_chart_to_png(chart, width_in=7.2, height_in=5.0, dpi=200)
            if img_buf:
                slide_c.shapes.add_picture(img_buf, Inches(0.8), Inches(1.4), width=Inches(7.2))

            # Right Side Context Box
            tb_right = slide_c.shapes.add_textbox(Inches(8.3), Inches(1.4), Inches(4.2), Inches(5.0))
            tf_r = tb_right.text_frame
            tf_r.word_wrap = True
            
            pr_head = tf_r.paragraphs[0]
            pr_head.text = "💡 Thông Điệp Số Liệu:"
            pr_head.font.bold = True
            pr_head.font.size = Pt(15)
            pr_head.font.color.rgb = TEXT_DARK

            # Top numbers
            labels = chart.get("labels") or []
            values = chart.get("values") or []
            for l, v in zip(labels[:6], values[:6]):
                pr_item = tf_r.add_paragraph()
                try:
                    pr_item.text = f"• {l}: {float(v):,.2f}"
                except Exception:
                    pr_item.text = f"• {l}: {v}"
                pr_item.font.size = Pt(12)
                pr_item.font.color.rgb = RGBColor(71, 85, 105)

    # ── FINAL SLIDE: Recommendations & Action Plan ──────────────────────
    recs = report.get("recommendations") or []
    if recs:
        slide_rec = prs.slides.add_slide(blank_layout)
        tb_rhdr = slide_rec.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        p_rh = tb_rhdr.text_frame.paragraphs[0]
        p_rh.text = "KẾ HOẠCH HÀNH ĐỘNG & ĐỀ XUẤT CHIẾN LƯỢC"
        p_rh.font.size = Pt(22)
        p_rh.font.bold = True
        p_rh.font.color.rgb = ACCENT_GREEN

        top_pos = 1.6
        for idx, rec in enumerate(recs[:4]):
            box = slide_rec.shapes.add_shape(1, Inches(0.8), Inches(top_pos), Inches(11.7), Inches(1.1))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(248, 250, 252)
            box.line.color.rgb = RGBColor(203, 213, 225)
            box.line.width = Pt(1)

            tf_box = box.text_frame
            tf_box.word_wrap = True
            p_step = tf_box.paragraphs[0]
            p_step.text = f"HÀNH ĐỘNG {idx+1}"
            p_step.font.bold = True
            p_step.font.size = Pt(10)
            p_step.font.color.rgb = ACCENT_GREEN

            p_desc = tf_box.add_paragraph()
            p_desc.text = str(rec)
            p_desc.font.size = Pt(12)
            p_desc.font.color.rgb = TEXT_DARK

            top_pos += 1.3

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3. Formatted Excel (.xlsx) Exporter with Chart Images & Formatted Tables
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def export_data_and_charts_to_xlsx(
    title: str,
    dataframes: dict[str, pd.DataFrame],
    charts: list[dict[str, Any]] | None = None,
    summary: str = "",
) -> io.BytesIO:
    """Render a clean, professional multi-sheet Excel workbook complete with embedded chart images."""
    import openpyxl
    from openpyxl.drawing.image import Image
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    
    # Sheet 1: Dashboard & Visualizations
    ws_dash = wb.active
    ws_dash.title = "📊 Báo Cáo & Biểu Đồ"
    ws_dash.views.sheetView[0].showGridLines = True

    # Styles
    font_title = Font(name="Segoe UI", size=16, bold=True, color="FFFFFF")
    fill_header = PatternFill(start_color="10B981", end_color="10B981", fill_type="solid")
    font_tbl_hdr = Font(name="Segoe UI", size=10, bold=True, color="FFFFFF")
    fill_tbl_hdr = PatternFill(start_color="1F2937", end_color="1F2937", fill_type="solid")
    fill_zebra = PatternFill(start_color="F9FAFB", end_color="F9FAFB", fill_type="solid")
    font_data = Font(name="Segoe UI", size=10)
    thin_border = Border(
        left=Side(style="thin", color="E5E7EB"), right=Side(style="thin", color="E5E7EB"),
        top=Side(style="thin", color="E5E7EB"), bottom=Side(style="thin", color="E5E7EB")
    )

    # Title Banner
    ws_dash.merge_cells("A1:K2")
    cell_top = ws_dash["A1"]
    cell_top.value = f"  {title.upper()}"
    cell_top.font = font_title
    cell_top.fill = fill_header
    cell_top.alignment = Alignment(vertical="center")

    current_row = 4
    if summary:
        ws_dash.cell(row=current_row, column=1, value="📌 Tóm tắt phân tích:").font = Font(name="Segoe UI", size=11, bold=True, color="10B981")
        current_row += 1
        ws_dash.merge_cells(start_row=current_row, start_column=1, end_row=current_row+1, end_column=11)
        sum_cell = ws_dash.cell(row=current_row, column=1, value=summary)
        sum_cell.font = Font(name="Segoe UI", size=10)
        sum_cell.alignment = Alignment(wrap_text=True, vertical="top")
        current_row += 3

    # Insert Charts as Images
    if charts:
        img_col = "A"
        for idx, chart in enumerate(charts[:3]):
            img_buf = render_chart_to_png(chart, width_in=6.5, height_in=3.8, dpi=180)
            if img_buf:
                try:
                    img = Image(img_buf)
                    img.anchor = f"A{current_row}"
                    ws_dash.add_image(img)
                    current_row += 20
                except Exception as exc:
                    logger.warning(f"[exporter] Could not add chart image to Excel: {exc}")

    # Subsequent Sheets: Formatted DataFrames
    for sheet_name, df in dataframes.items():
        if df is None or len(df) == 0:
            continue
        
        clean_sheet_name = str(sheet_name)[:30].replace("/", "_").replace("\\", "_")
        ws = wb.create_sheet(title=f"📄 {clean_sheet_name}")
        ws.views.sheetView[0].showGridLines = True

        # Header row
        for col_num, col_name in enumerate(df.columns, 1):
            cell = ws.cell(row=1, column=col_num, value=str(col_name))
            cell.font = font_tbl_hdr
            cell.fill = fill_tbl_hdr
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

        # Data rows with formatting
        for r_idx, row in enumerate(df.itertuples(index=False), 2):
            is_even = (r_idx % 2 == 0)
            for c_idx, val in enumerate(row, 1):
                cell = ws.cell(row=r_idx, column=c_idx)
                
                # Format numbers vs text
                if isinstance(val, (int, float)) and not pd.isna(val):
                    cell.value = val
                    if isinstance(val, float):
                        cell.number_format = "#,##0.00"
                    else:
                        cell.number_format = "#,##0"
                    cell.alignment = Alignment(horizontal="right")
                else:
                    cell.value = str(val) if not pd.isna(val) else ""
                    cell.alignment = Alignment(horizontal="left")

                cell.font = font_data
                cell.border = thin_border
                if is_even:
                    cell.fill = fill_zebra

        # Auto-adjust column widths
        for col in ws.columns:
            max_len = max(len(str(cell.value or "")) for cell in col)
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = min(max(max_len + 3, 12), 45)

    out = io.BytesIO()
    wb.save(out)
    out.seek(0)
    return out
